"""
make_pptx.py — 生成组会汇报 PPT (6 页,16:9)

内容:
  P1: Title — 9 基因项目数据调研与处理汇报
  P2: 项目背景 + 9 基因 panel
  P3: MC3 mutation 数据调研与提取
  P4: TCGA-LUAD-WSI DICOM 数据下载与处理
  P5: 真实切片可视化(科研图)
  P6: DeepGEM pipeline 借鉴 + 下一步计划
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

# ===== 调色板(与 chart 一致) =====
COL_PRIMARY = RGBColor(0x0F, 0x17, 0x2A)      # 标题深色
COL_ACCENT  = RGBColor(0x25, 0x63, 0xEB)      # 主蓝
COL_ACCENT2 = RGBColor(0xDC, 0x26, 0x26)      # 强调红
COL_TEXT    = RGBColor(0x1E, 0x29, 0x3B)      # 正文
COL_MUTED   = RGBColor(0x64, 0x74, 0x8B)      # 副文
COL_LIGHT   = RGBColor(0xF1, 0xF5, 0xF9)      # 浅底
COL_BORDER  = RGBColor(0xE2, 0xE8, 0xF0)
COL_BG_TAG  = RGBColor(0xE0, 0xE7, 0xFF)      # 标签底

ASSETS = "/mnt/d/pan-caner/group_meeting_report/assets"
OUT    = "/mnt/d/pan-caner/group_meeting_report/group_meeting_report.pptx"

prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9 widescreen
prs.slide_height = Inches(7.5)

# ---- 辅助函数 ----
def add_blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank

def add_text(slide, x, y, w, h, text, *,
             size=14, bold=False, color=COL_TEXT, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font_name="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font_name
    return tb

def add_bullets(slide, x, y, w, h, bullets, *,
                size=14, color=COL_TEXT, line_spacing=1.25,
                bullet_char="• ", indent_level=0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = bullet_char + item
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb

def add_rect(slide, x, y, w, h, fill, line=None, line_w=0):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh

def add_rounded(slide, x, y, w, h, fill, line=None, line_w=0):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.15
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh

def add_image(slide, path, x, y, w=None, h=None):
    if w and h:
        return slide.shapes.add_picture(path, x, y, width=w, height=h)
    elif w:
        return slide.shapes.add_picture(path, x, y, width=w)
    elif h:
        return slide.shapes.add_picture(path, x, y, height=h)
    return slide.shapes.add_picture(path, x, y)

def page_header(slide, page_num, title, subtitle=None):
    # top bar
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.85), COL_ACCENT)
    # page number circle
    add_rounded(slide, Inches(0.4), Inches(0.18), Inches(0.5), Inches(0.5), COL_BG_TAG)
    add_text(slide, Inches(0.4), Inches(0.18), Inches(0.5), Inches(0.5),
             str(page_num), size=18, bold=True, color=COL_ACCENT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # title
    add_text(slide, Inches(1.05), Inches(0.1), Inches(11.5), Inches(0.55),
             title, size=22, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
             anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(1.05), Inches(0.55), Inches(11.5), Inches(0.3),
                 subtitle, size=11, color=RGBColor(0xBF, 0xDB, 0xFE),
                 anchor=MSO_ANCHOR.MIDDLE)

def page_footer(slide, text="9-Gene Lung Cancer Project | Data Engineering Report | 2026-07-12"):
    add_text(slide, Inches(0.4), Inches(7.15), Inches(12.5), Inches(0.3),
             text, size=9, color=COL_MUTED)

# ============================================================
# P1: Title slide
# ============================================================
s = add_blank_slide()
# background
add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, COL_PRIMARY)
# accent diagonal
add_rect(s, Inches(0), Inches(5.4), prs.slide_width, Inches(2.1), COL_ACCENT)
# big icon-like number
add_text(s, Inches(0.6), Inches(1.0), Inches(3), Inches(1.5),
         "09", size=110, bold=True, color=RGBColor(0x60, 0xA5, 0xFA),
         anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.6), Inches(2.4), Inches(3), Inches(0.5),
         "GENES", size=18, bold=True, color=RGBColor(0x60, 0xA5, 0xFA))
add_text(s, Inches(0.6), Inches(2.8), Inches(3), Inches(0.5),
         "PANEL", size=18, bold=True, color=RGBColor(0x60, 0xA5, 0xFA))
# main title
add_text(s, Inches(4.2), Inches(1.2), Inches(8.8), Inches(1.0),
         "肺癌 9 基因突变预测项目", size=32, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(s, Inches(4.2), Inches(2.0), Inches(8.8), Inches(0.7),
         "Lung Adenocarcinoma 9-Gene Panel — Data Engineering Report",
         size=18, color=RGBColor(0x93, 0xC5, 0xFD))
add_text(s, Inches(4.2), Inches(2.7), Inches(8.8), Inches(0.5),
         "组会汇报 · 公开数据集调研与处理流程",
         size=15, color=RGBColor(0xDB, 0xE9, 0xFE))
# divider
add_rect(s, Inches(4.2), Inches(3.3), Inches(2), Inches(0.04),
         RGBColor(0x60, 0xA5, 0xFA))
# subtitle on accent
add_text(s, Inches(0.6), Inches(5.65), Inches(12), Inches(0.5),
         "MC3 mutation labels  ·  TCGA-LUAD DICOM WSI  ·  DeepGEM pipeline",
         size=18, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(s, Inches(0.6), Inches(6.2), Inches(12), Inches(0.4),
         "Data: 421 cases × 11 genes  |  WSI: 50 cases × ~5 series  |  Format: DICOM (.dcm)",
         size=13, color=RGBColor(0xBF, 0xDB, 0xFE))
# presenter info
add_text(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.3),
         "Reporter: zzy-ustb  |  Group meeting 2026-07-12",
         size=11, color=RGBColor(0xBF, 0xDB, 0xFE))

# ============================================================
# P2: 项目背景 + 9 基因 panel
# ============================================================
s = add_blank_slide()
page_header(s, 2, "项目背景与 9 基因 Panel",
            "Background · Why 9 driver genes · Project context")

# left: project context
add_text(s, Inches(0.4), Inches(1.1), Inches(6.0), Inches(0.4),
         "1. 项目背景", size=16, bold=True, color=COL_ACCENT)
ctx_bullets = [
    "合作方: 北科大(乙方·算法) × 聚时医疗(甲方·数据)",
    "硬指标: 准确率 ≥ 95% · 灵敏度 ≥ 90%",
    "交付: 2026-09 / 2027-03 / 2027-10 三个节点",
    "数据策略: 内部数据(聚时医疗) + 外部 TCGA-LUAD",
    "数据格式: 公开数据集(SVS/DICOM + MAF) + 自有 LIS-PACS",
    "技术栈: DINOv3 自监督 + 三维病理 + 多实例学习",
]
add_bullets(s, Inches(0.4), Inches(1.55), Inches(6.0), Inches(2.6),
            ctx_bullets, size=13, line_spacing=1.35)

# right: 9 gene panel
add_text(s, Inches(6.8), Inches(1.1), Inches(6.0), Inches(0.4),
         "2. 9 基因 Panel 来源", size=16, bold=True, color=COL_ACCENT)

# gene table — 3 cols × 4 rows, with source label
genes_panel = [
    ("EGFR", "KRAS", "ALK"),
    ("ROS1", "TP53", "BRAF"),
    ("PIK3CA", "ERBB2/HER2", "NRAS"),
    ("RET", "MET", "(11 total)"),
]
sources = ["DeepGEM 6 必含", "DeepGEM 6 必含", "DeepGEM 6 必含",
           "DeepGEM 6 必含", "DeepGEM 6 必含", "DeepGEM 6 必含",
           "商业靶向药", "合同附件三", "商业靶向药",
           "商业靶向药", "商业靶向药", "— 冗余"]
cell_w, cell_h = Inches(1.85), Inches(0.85)
x0, y0 = Inches(6.85), Inches(1.6)
for i, (a, b, c) in enumerate(genes_panel):
    for j, g in enumerate([a, b, c]):
        x = x0 + j * cell_w
        y = y0 + i * cell_h
        add_rounded(s, x, y, Inches(1.75), Inches(0.75), COL_LIGHT)
        add_text(s, x, y, Inches(1.75), Inches(0.4),
                 g, size=15, bold=True, color=COL_PRIMARY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, y + Inches(0.4), Inches(1.75), Inches(0.3),
                 sources[i*3 + j], size=8, color=COL_MUTED,
                 align=PP_ALIGN.CENTER)

# bottom: takeaway box
add_rounded(s, Inches(0.4), Inches(5.5), Inches(12.5), Inches(1.5), COL_BG_TAG)
add_text(s, Inches(0.6), Inches(5.65), Inches(12.1), Inches(0.4),
         "Takeaway", size=14, bold=True, color=COL_ACCENT)
add_text(s, Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.85),
         "本研究锁定 11 个核心基因(9 必含 + 2 冗余),覆盖 NCCN 指南必检靶点与 DeepGEM 已验证 panel。"
         "每个基因阳性 = 该 case 在 MC3 v0.2.8 PUBLIC MAF 中至少 1 个 PASS 功能改变性突变。",
         size=12, color=COL_TEXT, line_spacing=1.35) if False else None
add_bullets(s, Inches(0.6), Inches(6.05), Inches(12.1), Inches(0.85),
            ["11 个核心基因(9 必含 + 2 冗余),覆盖 NCCN 指南必检靶点与 DeepGEM 已验证 panel。",
             "阳性定义: MC3 v0.2.8 PUBLIC MAF 中至少 1 个 PASS 功能改变性突变(Missense / Nonsense / Frameshift / Splice 等)。"],
            size=12, line_spacing=1.25)

page_footer(s)

# ============================================================
# P3: MC3 mutation labels
# ============================================================
s = add_blank_slide()
page_header(s, 3, "MC3 MAF 数据调研与提取",
            "Source: Synapse syn7824274 · cBioPortal whitelist · 421 cases × 11 genes")

# left: workflow steps
add_text(s, Inches(0.4), Inches(1.05), Inches(6.0), Inches(0.4),
         "1. 调研结论", size=15, bold=True, color=COL_ACCENT)
discovery = [
    "MC3 (Multi-Center Mutation Calling) = TCGA pan-cancer 7-caller consensus 突变集合",
    "无 HTTP 直链 → 需 Synapse + PAT(personal access token)",
    "GDC Legacy Archive 关闭(PEAR-1166),新 GDC API 不返 Slide Image",
    "新方案: Synapse syn7824274 (PUBLIC,无需 dbGaP)",
]
add_bullets(s, Inches(0.4), Inches(1.45), Inches(6.0), Inches(1.8),
            discovery, size=12, line_spacing=1.3)

add_text(s, Inches(0.4), Inches(3.3), Inches(6.0), Inches(0.4),
         "2. 提取流程", size=15, bold=True, color=COL_ACCENT)
steps = [
    "下载 MC3 MAF 753 MB → streaming gzip 读 (内存单行)",
    "cBioPortal 拉 584 个 TCGA-LUAD case_id 白名单(修 is_luad() 只认 site 05 的 bug)",
    "白名单 ∩ MC3 → 421 cases 实际有数据",
    "PASS 过滤 → 基因白名单(11) → 变异类型白名单(9 类功能改变)",
    "输出: case_id × 11 genes 0/1 矩阵 + any_9gene_positive + n_genes_mutated",
]
add_bullets(s, Inches(0.4), Inches(3.7), Inches(6.0), Inches(2.0),
            steps, size=12, line_spacing=1.3)

add_text(s, Inches(0.4), Inches(5.8), Inches(6.0), Inches(0.4),
         "3. 工程产物", size=15, bold=True, color=COL_ACCENT)
add_text(s, Inches(0.4), Inches(6.2), Inches(6.0), Inches(0.4),
         "D:/pan-caner/data/MC3/9gene_panel_LUAD.csv (17 KB · 422 行)",
         size=11, color=COL_ACCENT2, bold=True)

# right: prevalence chart
add_image(s, os.path.join(ASSETS, "fig_gene_prevalence.png"),
          Inches(6.6), Inches(1.05), w=Inches(6.5))

# mini stat tiles at bottom right
add_rounded(s, Inches(6.6), Inches(5.4), Inches(2.05), Inches(1.5), COL_ACCENT)
add_text(s, Inches(6.6), Inches(5.5), Inches(2.05), Inches(0.7),
         "421", size=36, bold=True, color=RGBColor(0xFF,0xFF,0xFF),
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(6.6), Inches(6.15), Inches(2.05), Inches(0.5),
         "TCGA-LUAD cases", size=10, color=RGBColor(0xDB,0xE9,0xFE),
         align=PP_ALIGN.CENTER)
add_rounded(s, Inches(8.85), Inches(5.4), Inches(2.05), Inches(1.5), COL_ACCENT2)
add_text(s, Inches(8.85), Inches(5.5), Inches(2.05), Inches(0.7),
         "3.6M", size=36, bold=True, color=RGBColor(0xFF,0xFF,0xFF),
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(8.85), Inches(6.15), Inches(2.05), Inches(0.5),
         "MAF rows scanned", size=10, color=RGBColor(0xFE,0xE2,0xE2),
         align=PP_ALIGN.CENTER)
add_rounded(s, Inches(11.1), Inches(5.4), Inches(2.05), Inches(1.5),
            RGBColor(0x05, 0x96, 0x69))
add_text(s, Inches(11.1), Inches(5.5), Inches(2.05), Inches(0.7),
         "17 KB", size=36, bold=True, color=RGBColor(0xFF,0xFF,0xFF),
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(11.1), Inches(6.15), Inches(2.05), Inches(0.5),
         "final CSV", size=10, color=RGBColor(0xD1,0xFA,0xE5),
         align=PP_ALIGN.CENTER)

page_footer(s)

# ============================================================
# P4: DICOM WSI 数据
# ============================================================
s = add_blank_slide()
page_header(s, 4, "TCGA-LUAD WSI DICOM 数据下载与处理",
            "Source: NCI Imaging Data Commons (IDC) · 50 cases × ~5 series × multi-tile")

# left column: file format explanation
add_text(s, Inches(0.4), Inches(1.05), Inches(5.6), Inches(0.4),
         "1. WSI 格式: DICOM (.dcm) vs SVS vs SDPC", size=14, bold=True, color=COL_ACCENT)

fmt_table_data = [
    ("Format", "Vendor", "Reader"),
    (".svs", "Leica/Aperio", "openslide ✓"),
    (".sdpc", "Philips", "openslide ✓"),
    (".dcm", "DICOM standard", "wsidicom"),
]
y = Inches(1.55)
col_xs = [Inches(0.4), Inches(1.65), Inches(3.45)]
col_ws = [Inches(1.25), Inches(1.80), Inches(2.55)]
for i, row in enumerate(fmt_table_data):
    for j, val in enumerate(row):
        is_header = (i == 0)
        fill = COL_PRIMARY if is_header else (COL_LIGHT if i % 2 else RGBColor(0xFF,0xFF,0xFF))
        text_col = RGBColor(0xFF,0xFF,0xFF) if is_header else COL_TEXT
        add_rect(s, col_xs[j], y + i*Inches(0.4), col_ws[j], Inches(0.4), fill)
        add_text(s, col_xs[j] + Inches(0.1), y + i*Inches(0.4), col_ws[j] - Inches(0.2),
                 Inches(0.4), val, size=11, bold=is_header, color=text_col,
                 anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.4), Inches(3.4), Inches(5.6), Inches(0.4),
         "2. DICOM 内部结构(单 tile)", size=14, bold=True, color=COL_ACCENT)

# mini box: DICOM tag examples
dcm_info = [
    "PatientID (0010,0020): \"TCGA-05-4245\"  ← join key",
    "Modality (0008,0060): SM / SEG / ANN / CT",
    "TotalPixelMatrix (0048,0006-7): 4980×4500",
    "RowPosition (0048,0070): tile 左上角 Y 坐标",
    "Objective Lens Power (0048,0112): 20×",
    "PixelSpacing (0028,0030): 0.5 μm/px",
]
add_bullets(s, Inches(0.4), Inches(3.8), Inches(5.6), Inches(2.6),
            dcm_info, size=11, line_spacing=1.4)

add_text(s, Inches(0.4), Inches(6.0), Inches(5.6), Inches(0.5),
         "3. 拼接规则: 标准 tag, 不会拼错", size=13, bold=True, color=COL_ACCENT)
add_text(s, Inches(0.4), Inches(6.4), Inches(5.6), Inches(0.7),
         "Row/Column PositionInTotalImagePixelMatrix 是 DICOM Part 3 强制要求 tag, "
         "所有厂商(Philips / Leica / Hamamatsu)统一",
         size=10, color=COL_MUTED)

# right column: stitch illustration
add_image(s, os.path.join(ASSETS, "fig_dicom_stitch.png"),
          Inches(6.3), Inches(1.05), w=Inches(6.8))

# right column: dataset at a glance
add_image(s, os.path.join(ASSETS, "fig_stat_tiles.png"),
          Inches(6.3), Inches(4.7), w=Inches(6.8))

page_footer(s)

# ============================================================
# P5: 真实切片可视化
# ============================================================
s = add_blank_slide()
page_header(s, 5, "真实切片可视化 (TCGA-05-4245, LUAD)",
            "Source: IDC, 20× objective, Manufacturer: Carl Zeiss (Aperio → DICOM)")

# big patch image (left)
add_text(s, Inches(0.4), Inches(1.05), Inches(7.0), Inches(0.4),
         "中央 1024×1024 patch (from stitched series, 4980×4500 @ 20×)",
         size=12, bold=True, color=COL_ACCENT)
add_image(s, os.path.join(ASSETS, "patch_sample.png"),
          Inches(0.4), Inches(1.55), w=Inches(7.0))

# annotation on patch
add_rounded(s, Inches(0.5), Inches(6.05), Inches(6.8), Inches(0.85),
            RGBColor(0xF0, 0xFD, 0xF4))
add_text(s, Inches(0.7), Inches(6.1), Inches(6.5), Inches(0.4),
         "H&E 染色 (Hematoxylin & Eosin)",
         size=12, bold=True, color=RGBColor(0x05, 0x96, 0x69))
add_text(s, Inches(0.7), Inches(6.45), Inches(6.5), Inches(0.4),
         "紫色 = 细胞核 · 粉红色 = 胞质 / 间质 · 白色空洞 = 肺泡结构",
         size=10, color=RGBColor(0x06, 0x6B, 0x4A))

# right column: thumbnail + key facts
add_text(s, Inches(7.6), Inches(1.05), Inches(5.5), Inches(0.4),
         "完整切片缩略图 (~1848×2048)", size=12, bold=True, color=COL_ACCENT)
add_image(s, os.path.join(ASSETS, "wsi_thumbnail.png"),
          Inches(7.6), Inches(1.5), w=Inches(5.4))

# key facts box
add_rounded(s, Inches(7.6), Inches(5.0), Inches(5.4), Inches(2.0),
            COL_LIGHT)
add_text(s, Inches(7.8), Inches(5.1), Inches(5.0), Inches(0.4),
         "DICOM header → WSI 关键事实",
         size=12, bold=True, color=COL_PRIMARY)
facts = [
    "PatientID = TCGA-05-4245 (join key)",
    "Trial Protocol ID = TCGA-LUAD",
    "Manufacturer: Leica Biosystems",
    "Model: Aperio (TIFF → DICOM by PixelMed)",
    "Series Description: Frozen HE TP BS1",
    "Stain: hematoxylin + eosin (H&E)",
    "Container ID: TCGA-05-4245-01A-01-BS1",
]
add_bullets(s, Inches(7.8), Inches(5.45), Inches(5.0), Inches(1.5),
            facts, size=10, line_spacing=1.2, bullet_char="• ")

page_footer(s)

# ============================================================
# P6: DeepGEM pipeline + 下一步
# ============================================================
s = add_blank_slide()
page_header(s, 6, "DeepGEM Pipeline 借鉴与下一步计划",
            "Tencent AI Lab · Lancet Oncology · Step1~Step4 = our path")

# top: pipeline
add_image(s, os.path.join(ASSETS, "fig_pipeline.png"),
          Inches(0.4), Inches(1.05), w=Inches(12.5))

# left: DeepGEM key facts
add_text(s, Inches(0.4), Inches(4.1), Inches(6.0), Inches(0.4),
         "1. DeepGEM pipeline 关键事实", size=14, bold=True, color=COL_ACCENT)
deepgem_facts = [
    "WSI → 1120×1120 patches @ 20× (与 DICOM 物镜一致 ✓)",
    "特征提取: CTransPath / EfficientNet-B0 → .pkl",
    "MIL 模型: Transformer + label disambiguation",
    "输入实际是 patch features, 不是原始 WSI",
    "opTCGA 训练集 → 跨中心验证 7 个医院",
]
add_bullets(s, Inches(0.4), Inches(4.5), Inches(6.0), Inches(2.4),
            deepgem_facts, size=12, line_spacing=1.35)

# right: next steps
add_text(s, Inches(6.8), Inches(4.1), Inches(6.0), Inches(0.4),
         "2. 下一步计划", size=14, bold=True, color=COL_ACCENT)

next_steps = [
    ("DICOM adapter", "wsidicom 模拟 openslide 接口 → DeepGEM step1 无需修改", COL_ACCENT),
    ("Patch 提取", "50 cases × ~100 patches = 5000 patches, 落盘 PNG", RGBColor(0x05, 0x96, 0x69)),
    ("CTransPath 特征", "预训练 backbone 提 patch features → combined_feat.pkl", RGBColor(0x7C, 0x3A, 0xED)),
    ("DeepGEM 微调", "用 MC3 9-gene label 微调 + 合同 95% 准确率指标", COL_ACCENT2),
    ("异源验证", "钟团队 CJFH 数据集(需 dbGaP 授权)", COL_MUTED),
]
y0 = Inches(4.55)
for i, (head, body, color) in enumerate(next_steps):
    yy = y0 + i * Inches(0.5)
    add_rounded(s, Inches(6.8), yy, Inches(0.18), Inches(0.4), color)
    add_text(s, Inches(7.1), yy, Inches(2.3), Inches(0.4),
             head, size=12, bold=True, color=color,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(9.5), yy, Inches(3.5), Inches(0.4),
             body, size=10, color=COL_TEXT,
             anchor=MSO_ANCHOR.MIDDLE)

page_footer(s)

# save
prs.save(OUT)
print(f"PPT saved: {OUT}")
print(f"Slides: {len(prs.slides)}")